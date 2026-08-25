from pathlib import Path

import pytest

from app.services.documents.attachment_parsers import AttachmentParserDispatch


def parse_file(tmp_path: Path, filename: str, content: str):
    path = tmp_path / filename
    path.write_text(content, encoding="utf-8")
    return AttachmentParserDispatch().parse(
        path,
        document_id="doc-1",
        original_filename=filename,
        max_pages=20,
        min_extracted_chars=1,
    )


@pytest.mark.parametrize(
    ("filename", "content_kind", "locator_unit"),
    [
        ("notes.txt", "text", "line"),
        ("script.py", "code", "line"),
        ("table.csv", "csv", "row"),
        ("data.json", "json", "section"),
        ("page.html", "html", "section"),
    ],
)
def test_dispatch_parses_textual_attachment_kinds(
    tmp_path,
    filename,
    content_kind,
    locator_unit,
):
    content = {
        "notes.txt": "one\ntwo\n",
        "script.py": "print('ok')\n",
        "table.csv": "name,value\nalpha,1\n",
        "data.json": '{"answer": 42}',
        "page.html": "<html><h1>Title</h1><p>Body</p></html>",
    }[filename]

    parsed = parse_file(tmp_path, filename, content)

    assert parsed.content_kind == content_kind
    assert parsed.locator_unit == locator_unit
    assert parsed.page_count == len(parsed.pages) > 0
    assert parsed.extracted_char_count > 0
    assert all(page.width == 0.0 and page.height == 0.0 for page in parsed.pages)
    assert all(block.bbox == (0.0, 0.0, 0.0, 0.0) for page in parsed.pages for block in page.blocks)


def test_text_is_split_into_120_line_virtual_pages(tmp_path):
    parsed = parse_file(
        tmp_path,
        "notes.txt",
        "\n".join(f"line {index}" for index in range(1, 241)),
    )

    assert parsed.page_count == 2
    assert (parsed.pages[0].locator_start, parsed.pages[0].locator_end) == (1, 120)
    assert (parsed.pages[1].locator_start, parsed.pages[1].locator_end) == (121, 240)


def test_markdown_heading_sections_get_section_locators(tmp_path):
    parsed = parse_file(
        tmp_path,
        "notes.md",
        "# Intro\nfirst\n## Details\nsecond\n",
    )

    assert parsed.content_kind == "markdown"
    assert parsed.locator_unit == "section"
    assert parsed.page_count == 2
    assert [page.locator for page in parsed.pages] == ["Intro", "Details"]


def test_parser_identity_is_extension_aware():
    dispatch = AttachmentParserDispatch()

    assert dispatch.identity_for("notes.pdf")[0] == "pymupdf"
    assert dispatch.identity_for("notes.md") == ("markdown", "1")
    assert dispatch.identity_for("notes.txt") == ("text", "1")
    assert dispatch.identity_for("notes.docx") == ("python-docx", "1")


def test_docx_uses_heading_sections_and_paragraph_locators(tmp_path):
    from docx import Document

    path = tmp_path / "notes.docx"
    document = Document()
    document.add_heading("Intro", level=1)
    document.add_paragraph("First section")
    document.add_heading("Details", level=2)
    document.add_paragraph("Second section")
    document.save(path)

    parsed = AttachmentParserDispatch().parse(
        path,
        document_id="doc-1",
        original_filename=path.name,
        max_pages=20,
        min_extracted_chars=1,
    )

    assert parsed.content_kind == "docx"
    assert parsed.locator_unit == "paragraph"
    assert [page.locator for page in parsed.pages] == ["Intro", "Details"]
    assert (parsed.pages[0].locator_start, parsed.pages[0].locator_end) == (1, 2)


def test_xlsx_creates_one_virtual_page_per_sheet(tmp_path):
    from openpyxl import Workbook

    path = tmp_path / "budget.xlsx"
    workbook = Workbook()
    summary = workbook.active
    summary.title = "Summary"
    summary.append(["Total", 42])
    detail = workbook.create_sheet("Detail")
    detail.append(["Item", "Cost"])
    detail.append(["Book", 20])
    workbook.save(path)
    workbook.close()

    parsed = AttachmentParserDispatch().parse(
        path,
        document_id="doc-1",
        original_filename=path.name,
        max_pages=20,
        min_extracted_chars=1,
    )

    assert parsed.content_kind == "xlsx"
    assert parsed.locator_unit == "sheet"
    assert [page.locator for page in parsed.pages] == ["Summary", "Detail"]
    assert "Total\t42" in parsed.pages[0].blocks[0].text


def test_pptx_creates_one_virtual_page_per_slide(tmp_path):
    from pptx import Presentation

    path = tmp_path / "lesson.pptx"
    presentation = Presentation()
    first = presentation.slides.add_slide(presentation.slide_layouts[1])
    first.shapes.title.text = "Intro"
    first.placeholders[1].text = "First slide"
    second = presentation.slides.add_slide(presentation.slide_layouts[1])
    second.shapes.title.text = "Details"
    second.placeholders[1].text = "Second slide"
    presentation.save(path)

    parsed = AttachmentParserDispatch().parse(
        path,
        document_id="doc-1",
        original_filename=path.name,
        max_pages=20,
        min_extracted_chars=1,
    )

    assert parsed.content_kind == "pptx"
    assert parsed.locator_unit == "slide"
    assert parsed.page_count == 2
    assert "Intro" in parsed.pages[0].blocks[0].text
    assert "Second slide" in parsed.pages[1].blocks[0].text
