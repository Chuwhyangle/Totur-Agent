"""Convert supported non-PDF attachments into the stable parsed schema."""

from pathlib import Path
import csv
import json
import re

from bs4 import BeautifulSoup

from app.services.documents.parsed_document import (
    ParsedDocument,
    ParsedPage,
    ParsedTextBlock,
)
from app.services.documents.pdf_parser import (
    NoExtractableText,
    PdfContentLimitExceeded,
    PdfParseFailed,
    PdfParser,
    PdfParsingError,
    PdfSourceUnavailable,
)


TEXT_EXTENSIONS = {
    ".txt", ".log", ".md", ".markdown", ".py", ".js", ".jsx", ".ts",
    ".tsx", ".java", ".go", ".rs", ".c", ".cpp", ".h", ".cs", ".sql",
    ".sh", ".yaml", ".yml", ".toml", ".ini", ".css", ".xml",
}


class AttachmentParserDispatch:
    """Route one trusted attachment path to the appropriate parser."""

    name = "pymupdf"
    version = PdfParser.version

    def __init__(self, pdf_parser: PdfParser | None = None) -> None:
        self.pdf_parser = pdf_parser or PdfParser()

    def identity_for(self, original_filename: str) -> tuple[str, str]:
        extension = Path(original_filename).suffix.lower()
        if extension == ".pdf":
            return self.pdf_parser.name, self.pdf_parser.version
        parser_name = {
            ".md": "markdown",
            ".markdown": "markdown",
            ".csv": "csv",
            ".json": "json",
            ".html": "html",
            ".docx": "python-docx",
            ".xlsx": "openpyxl",
            ".pptx": "python-pptx",
        }.get(extension, "text")
        return parser_name, "1"

    def parse(
        self,
        source_path: Path,
        document_id: str,
        original_filename: str,
        max_pages: int,
        min_extracted_chars: int,
        max_extracted_chars: int = 2_000_000,
        max_blocks_per_page: int = 5_000,
    ) -> ParsedDocument:
        extension = Path(original_filename).suffix.lower()
        if extension == ".pdf":
            return self.pdf_parser.parse(
                source_path=source_path,
                document_id=document_id,
                original_filename=original_filename,
                max_pages=max_pages,
                min_extracted_chars=min_extracted_chars,
                max_extracted_chars=max_extracted_chars,
                max_blocks_per_page=max_blocks_per_page,
            )

        try:
            if extension in {".md", ".markdown"}:
                pages, kind, unit = _parse_markdown(Path(source_path))
            elif extension in TEXT_EXTENSIONS:
                pages, kind, unit = _parse_line_text(Path(source_path))
            elif extension == ".csv":
                pages, kind, unit = _parse_csv(Path(source_path))
            elif extension == ".json":
                pages, kind, unit = _parse_json(Path(source_path))
            elif extension == ".html":
                pages, kind, unit = _parse_html(Path(source_path))
            elif extension == ".docx":
                pages, kind, unit = _parse_docx(Path(source_path))
            elif extension == ".xlsx":
                pages, kind, unit = _parse_xlsx(Path(source_path))
            elif extension == ".pptx":
                pages, kind, unit = _parse_pptx(Path(source_path))
            else:
                raise PdfParseFailed(
                    f"No parser is registered for {extension or 'this attachment'}"
                )
        except PdfParsingError:
            raise
        except (FileNotFoundError, PermissionError, OSError) as exc:
            raise PdfSourceUnavailable("Attachment source is unavailable") from exc
        except Exception as exc:
            raise PdfParseFailed("Attachment could not be parsed") from exc

        if len(pages) > max_pages:
            raise PdfContentLimitExceeded(
                f"Attachment has {len(pages)} sections; limit is {max_pages}"
            )
        extracted_chars = sum(
            sum(not character.isspace() for character in block.text)
            for page in pages
            for block in page.blocks
        )
        if extracted_chars > max_extracted_chars:
            raise PdfContentLimitExceeded(
                "Attachment extracted text exceeds the configured limit"
            )
        if extracted_chars < min_extracted_chars:
            raise NoExtractableText("Attachment does not contain enough text")

        return ParsedDocument(
            schema_version=2,
            document_id=document_id,
            original_filename=original_filename,
            page_count=len(pages),
            extracted_char_count=extracted_chars,
            pages=tuple(pages),
            content_kind=kind,
            locator_unit=unit,
        )


def _parse_line_text(path: Path) -> tuple[list[ParsedPage], str, str]:
    text = path.read_text(encoding="utf-8")
    lines = text.splitlines()
    pages = _pages_from_line_groups(lines, 120)
    extension = path.suffix.lower()
    kind = "log" if extension == ".log" else "text" if extension == ".txt" else "code"
    return pages, kind, "line"


def _parse_markdown(path: Path) -> tuple[list[ParsedPage], str, str]:
    lines = path.read_text(encoding="utf-8").splitlines()
    headings = [
        index
        for index, line in enumerate(lines)
        if re.match(r"^\s{0,3}#{1,2}\s+", line)
    ]
    if not headings:
        return _pages_from_line_groups(lines, 120), "markdown", "line"

    pages: list[ParsedPage] = []
    section_starts = ([0] if headings[0] > 0 else []) + headings
    boundaries = section_starts + [len(lines)]
    for section_index, start in enumerate(section_starts):
        end = boundaries[section_index + 1]
        text = "\n".join(lines[start:end]).strip()
        if text:
            heading = (
                lines[start].lstrip("# ").strip()
                if start in headings
                else None
            )
            pages.append(
                _page(
                    len(pages) + 1,
                    text,
                    locator_start=start + 1,
                    locator_end=end,
                    locator=heading,
                )
            )
    return pages, "markdown", "section"


def _parse_csv(path: Path) -> tuple[list[ParsedPage], str, str]:
    with path.open("r", encoding="utf-8", newline="") as source:
        rows = [", ".join(row) for row in csv.reader(source)]
    pages: list[ParsedPage] = []
    for start in range(0, len(rows), 200):
        group = rows[start : start + 200]
        text = "\n".join(group).strip()
        if text:
            pages.append(
                _page(
                    len(pages) + 1,
                    text,
                    locator_start=start + 1,
                    locator_end=start + len(group),
                )
            )
    return pages, "csv", "row"


def _parse_json(path: Path) -> tuple[list[ParsedPage], str, str]:
    value = json.loads(path.read_text(encoding="utf-8"))
    pretty = json.dumps(value, ensure_ascii=False, indent=2)
    pages = [
        _page(index + 1, pretty[start : start + 3000])
        for index, start in enumerate(range(0, len(pretty), 3000))
    ]
    return pages, "json", "section"


def _parse_html(path: Path) -> tuple[list[ParsedPage], str, str]:
    soup = BeautifulSoup(path.read_text(encoding="utf-8"), "html.parser")
    pages: list[ParsedPage] = []
    current: list[str] = []
    locator: str | None = None
    for element in soup.find_all(["h1", "h2", "p", "li", "pre", "td", "th"]):
        value = element.get_text(" ", strip=True)
        if not value:
            continue
        if element.name in {"h1", "h2"} and current:
            pages.append(_page(len(pages) + 1, "\n".join(current), locator=locator))
            current = []
        if element.name in {"h1", "h2"}:
            locator = value
        current.append(value)
    if current:
        pages.append(_page(len(pages) + 1, "\n".join(current), locator=locator))
    if not pages:
        text = soup.get_text("\n", strip=True)
        pages = _pages_from_line_groups(text.splitlines(), 120)
    return pages, "html", "section"


def _parse_docx(path: Path) -> tuple[list[ParsedPage], str, str]:
    from docx import Document

    paragraphs = list(Document(path).paragraphs)
    heading_indices = [
        index
        for index, paragraph in enumerate(paragraphs)
        if paragraph.style is not None
        and paragraph.style.name.lower() in {"heading 1", "heading 2"}
    ]
    starts = (
        (([0] if heading_indices[0] > 0 else []) + heading_indices)
        if heading_indices
        else list(range(0, len(paragraphs), 40))
    )
    boundaries = starts + [len(paragraphs)]
    pages: list[ParsedPage] = []
    for section_index, start in enumerate(starts):
        end = boundaries[section_index + 1]
        text = "\n".join(
            paragraph.text.strip()
            for paragraph in paragraphs[start:end]
            if paragraph.text.strip()
        )
        if not text:
            continue
        locator = (
            paragraphs[start].text.strip()
            if start in heading_indices
            else None
        )
        pages.append(
            _page(
                len(pages) + 1,
                text,
                locator_start=start + 1,
                locator_end=end,
                locator=locator,
            )
        )
    return pages, "docx", "paragraph"


def _parse_xlsx(path: Path) -> tuple[list[ParsedPage], str, str]:
    from openpyxl import load_workbook

    workbook = load_workbook(path, read_only=True, data_only=True)
    try:
        pages = []
        for sheet in workbook.worksheets:
            rows = []
            for row in sheet.iter_rows(values_only=True):
                values = ["" if value is None else str(value) for value in row]
                if any(values):
                    rows.append("\t".join(values).rstrip())
            pages.append(
                _page(
                    len(pages) + 1,
                    "\n".join(rows),
                    locator=sheet.title,
                )
            )
    finally:
        workbook.close()
    return pages, "xlsx", "sheet"


def _parse_pptx(path: Path) -> tuple[list[ParsedPage], str, str]:
    from pptx import Presentation

    presentation = Presentation(path)
    pages: list[ParsedPage] = []
    for slide in presentation.slides:
        texts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                value = shape.text.strip()
                if value:
                    texts.append(value)
            elif getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    value = "\t".join(cell.text.strip() for cell in row.cells).rstrip()
                    if value:
                        texts.append(value)
        pages.append(_page(len(pages) + 1, "\n".join(texts)))
    return pages, "pptx", "slide"


def _pages_from_line_groups(lines: list[str], group_size: int) -> list[ParsedPage]:
    pages: list[ParsedPage] = []
    for start in range(0, len(lines), group_size):
        group = lines[start : start + group_size]
        text = "\n".join(group).strip()
        if text:
            pages.append(
                _page(
                    len(pages) + 1,
                    text,
                    locator_start=start + 1,
                    locator_end=start + len(group),
                )
            )
    return pages


def _page(
    page_number: int,
    text: str,
    *,
    locator_start: int | None = None,
    locator_end: int | None = None,
    locator: str | None = None,
) -> ParsedPage:
    normalized = text.strip()
    blocks = (
        (
            ParsedTextBlock(
                block_index=0,
                text=normalized,
                bbox=(0.0, 0.0, 0.0, 0.0),
            ),
        )
        if normalized
        else ()
    )
    return ParsedPage(
        page_number=page_number,
        width=0.0,
        height=0.0,
        blocks=blocks,
        locator_start=locator_start,
        locator_end=locator_end,
        locator=locator,
    )
